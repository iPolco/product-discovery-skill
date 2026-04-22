#!/usr/bin/env python3
"""Вставить слайд «Сравнение конкурентов» после слайда 7 (Конкурентный ландшафт).

Новый слайд содержит заготовку таблицы 6×N со столбцами:
Конкурент | Рынок | Цена / Подписка | Клиенты (оценка) | Выручка (оценка) | Ключевые фичи | Слабость

Запускается один раз на шаблоне — после этого слайд становится частью базовой структуры
презентации. Если шаблон уже был обновлён, скрипт выведет предупреждение и ничего не сделает.

Usage:
    python3 add_competitor_comparison_slide.py <path-to-presentation.pptx>

Файл модифицируется на месте.
"""
import copy
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("❌ Нужна библиотека python-pptx. Установить: pip install python-pptx")
    sys.exit(1)


SLIDE_TITLE = "Сравнение конкурентов: финансы и продукт"
MARKER_TEXT = "Сравнение конкурентов"  # используется для проверки, есть ли уже слайд

HEADERS = [
    "Конкурент",
    "Рынок / гео",
    "Цена / подписка",
    "Клиенты (оценка)",
    "Выручка (оценка)",
    "Ключевые фичи",
    "Слабость",
]


def has_comparison_slide(prs) -> bool:
    """Проверяет, есть ли уже слайд сравнения конкурентов."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and MARKER_TEXT in shape.text_frame.text:
                return True
    return False


def move_slide(prs, old_index: int, new_index: int):
    """Перемещает слайд с old_index на new_index в порядке презентации."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    slide_to_move = slides_list[old_index]
    xml_slides.remove(slide_to_move)
    xml_slides.insert(new_index, slide_to_move)


def add_competitor_comparison_slide(pptx_path: str) -> bool:
    """Вставить слайд после 7-го. Возвращает True если успешно, False если уже есть."""
    prs = Presentation(pptx_path)

    if has_comparison_slide(prs):
        print("ℹ️  Слайд сравнения конкурентов уже есть в презентации, пропускаю.")
        return False

    layout = prs.slide_layouts[0]  # DEFAULT
    # add_slide добавляет в конец; потом переставим
    slide = prs.slides.add_slide(layout)
    new_slide_index = len(prs.slides) - 1

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Заголовок слайда
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.25), slide_w - Inches(0.8), Inches(0.6)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = SLIDE_TITLE
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # Подзаголовок / инструкция агенту
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.85), slide_w - Inches(0.8), Inches(0.4)
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    sp = subtitle_tf.paragraphs[0]
    srun = sp.add_run()
    srun.text = (
        "Финансовые и продуктовые параметры основных игроков. "
        "Цифры — оценка по открытым источникам, указывай уровень достоверности."
    )
    srun.font.name = "Arial"
    srun.font.size = Pt(11)
    srun.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
    srun.font.italic = True

    # Таблица 6 строк (1 заголовок + 5 конкурентов) × 7 колонок
    rows_count = 6
    cols_count = len(HEADERS)

    table_left = Inches(0.4)
    table_top = Inches(1.4)
    table_width = slide_w - Inches(0.8)
    table_height = Inches(3.8)

    table_shape = slide.shapes.add_table(rows_count, cols_count, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # Распределение ширины колонок: Конкурент и Ключевые фичи шире, остальные равномерно
    col_weights = [1.2, 0.9, 1.0, 0.9, 0.9, 1.5, 1.0]  # относительные веса
    total_weight = sum(col_weights)
    for i, w in enumerate(col_weights):
        table.columns[i].width = Emu(int((slide_w - Inches(0.8)) * w / total_weight))

    # Заголовки
    for i, header in enumerate(HEADERS):
        cell = table.cell(0, i)
        cell.text = header
        # Стилизуем заголовок
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Фон заголовка
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # Строки-плейсхолдеры
    for row_idx in range(1, rows_count):
        for col_idx in range(cols_count):
            cell = table.cell(row_idx, col_idx)
            placeholder = f"[Конкурент {row_idx}]" if col_idx == 0 else "[...]"
            cell.text = placeholder
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if col_idx in (0, 5) else PP_ALIGN.CENTER
            for run in para.runs:
                run.font.name = "Arial"
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    # Подпись источника
    src_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(5.3), slide_w - Inches(0.8), Inches(0.3)
    )
    src_tf = src_box.text_frame
    src_p = src_tf.paragraphs[0]
    src_run = src_p.add_run()
    src_run.text = "Источники: G2, Crunchbase, SimilarWeb, отчёты компаний (задача 3, 4)"
    src_run.font.name = "Arial"
    src_run.font.size = Pt(8)
    src_run.font.italic = True
    src_run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # Переставляем созданный слайд сразу после слайда 7 (индекс 6)
    # Цель: новый слайд на индексе 7 (т.е. 8-й по порядку)
    target_index = 7
    move_slide(prs, new_slide_index, target_index)

    prs.save(pptx_path)

    # ══════════════════════════════════════════════════════════════════
    # КРИТИЧНО: чистка дублированных ссылок на notesSlides
    # ══════════════════════════════════════════════════════════════════
    # python-pptx при копировании слайдов копирует и _rels, что приводит
    # к тому что несколько слайдов ссылаются на один notesSlide. PowerPoint
    # из-за этого отказывается открывать файл (хотя LibreOffice — открывает).
    # См. references/block-6-artifacts.md раздел «Troubleshooting pack.py».
    cleanup_duplicate_notes_refs(pptx_path)

    print(f"✅ Слайд «{SLIDE_TITLE}» вставлен на позицию {target_index + 1}.")
    print(f"   Теперь в презентации {len(prs.slides)} слайдов.")
    print("   ⚠️  Проверь, что нумерация слайдов в SKILL.md / block-6-artifacts.md обновлена.")
    return True


def cleanup_duplicate_notes_refs(pptx_path):
    """Удаляет дублирующие ссылки на notesSlides, возникающие после копирования слайдов.

    PowerPoint требует уникальные notes на каждый slide. При копировании rels через
    python-pptx новый slide ссылается на тот же notesSlide — это ломает файл в
    PowerPoint (валидация по OOXML проходит, но PowerPoint её не пропускает).

    Решение — удалить все ссылки на notesSlides из slide rels и саму папку
    notesSlides. Заметки для презентаций инвестору не нужны.
    """
    import zipfile
    import tempfile
    import shutil
    import re
    import os

    tmp = tempfile.mkdtemp()
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            z.extractall(tmp)

        # 1. Удалить ссылки из slide.xml.rels
        rels_dir = os.path.join(tmp, "ppt/slides/_rels")
        if os.path.exists(rels_dir):
            for fname in os.listdir(rels_dir):
                path = os.path.join(rels_dir, fname)
                with open(path, encoding="utf-8") as fp:
                    content = fp.read()
                # НЕ-жадный матч: Target содержит '/' в пути, поэтому [^>]*?
                new = re.sub(r'\s*<Relationship[^>]*?notesSlide[^>]*?/>', '', content)
                if new != content:
                    with open(path, "w", encoding="utf-8") as fp:
                        fp.write(new)

        # 2. Удалить саму папку notesSlides
        notes_dir = os.path.join(tmp, "ppt/notesSlides")
        if os.path.exists(notes_dir):
            shutil.rmtree(notes_dir)

        # 3. Удалить Override для notesSlides из [Content_Types].xml
        ct_path = os.path.join(tmp, "[Content_Types].xml")
        if os.path.exists(ct_path):
            with open(ct_path, encoding="utf-8") as fp:
                ct = fp.read()
            new_ct = re.sub(r'\s*<Override[^>]*?notesSlide[^>]*?/>', '', ct)
            if new_ct != ct:
                with open(ct_path, "w", encoding="utf-8") as fp:
                    fp.write(new_ct)

        # 4. Пересобрать zip
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as z:
            for root, _, files in os.walk(tmp):
                for fname in files:
                    fp = os.path.join(root, fname)
                    z.write(fp, os.path.relpath(fp, tmp))
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(__doc__)
        sys.exit(1)

    path = Path(sys.argv[1])
    if not path.exists():
        print(f"❌ Файл не найден: {path}")
        sys.exit(1)

    add_competitor_comparison_slide(str(path))
